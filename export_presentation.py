#!/usr/bin/env python3
"""Export phase1.html slides to PowerPoint and PDF."""
import os
import time
from pathlib import Path
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from PIL import Image

BASE = Path(__file__).parent
HTML = BASE / "presentation" / "phase1.html"
OUT_DIR = BASE / "presentation"
SLIDE_COUNT = 7
WIDTH, HEIGHT = 1920, 1080

def take_screenshots():
    imgs = []
    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport={"width": WIDTH, "height": HEIGHT})
        page.goto(f"file://{HTML}")
        # wait for fonts/icons
        page.wait_for_timeout(2000)

        for i in range(SLIDE_COUNT):
            path = OUT_DIR / f"slide_{i+1:02d}.png"
            page.screenshot(path=str(path), clip={"x": 0, "y": 0, "width": WIDTH, "height": HEIGHT})
            imgs.append(path)
            print(f"  Screenshot Slide {i+1}/{SLIDE_COUNT}")
            if i < SLIDE_COUNT - 1:
                # click next
                page.click("#btn-next")
                page.wait_for_timeout(400)

        browser.close()
    return imgs

def build_pptx(imgs):
    prs = Presentation()
    prs.slide_width  = Emu(WIDTH * 9144)   # 1920 px @ 96 dpi → EMU
    prs.slide_height = Emu(HEIGHT * 9144)
    # 9144 EMU per pixel at 96 dpi

    blank_layout = prs.slide_layouts[6]  # completely blank

    for path in imgs:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(path), 0, 0,
            width=prs.slide_width,
            height=prs.slide_height
        )
        print(f"  PPTX: {path.name}")

    out = OUT_DIR / "phase1_presentation.pptx"
    prs.save(str(out))
    return out

def build_pdf(imgs):
    from PIL import Image as PILImage
    out = OUT_DIR / "phase1_presentation.pdf"

    pil_imgs = [PILImage.open(str(p)).convert("RGB") for p in imgs]
    first = pil_imgs[0]
    rest  = pil_imgs[1:]
    first.save(
        str(out),
        save_all=True,
        append_images=rest,
        resolution=150
    )
    return out

if __name__ == "__main__":
    print("=== Starte Export ===")
    print("1) Screenshots …")
    imgs = take_screenshots()

    print("2) PowerPoint erstellen …")
    pptx_out = build_pptx(imgs)
    print(f"   → {pptx_out}")

    print("3) PDF erstellen …")
    pdf_out = build_pdf(imgs)
    print(f"   → {pdf_out}")

    # cleanup screenshots
    print("4) Aufräumen …")
    for p in imgs:
        p.unlink()

    print("=== Fertig ===")
    print(f"  PPTX: {pptx_out}")
    print(f"  PDF:  {pdf_out}")
