"""Export index.html → pitch.pptx (pixel-perfect screenshots via Playwright)."""

import asyncio
from pathlib import Path
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches

SLIDE_W = 1920
SLIDE_H = 1080
HTML_PATH = Path(__file__).parent / "index.html"
OUT_PATH  = Path(__file__).parent / "pitch.pptx"


async def capture_slides():
    screenshots = []
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page = await browser.new_page(viewport={"width": SLIDE_W, "height": SLIDE_H})
        await page.goto(f"file://{HTML_PATH.resolve()}")
        await page.wait_for_load_state("networkidle")

        n_slides = await page.evaluate("document.querySelectorAll('.slide').length")
        print(f"Found {n_slides} slides")

        for i in range(n_slides):
            await page.evaluate(f"""
                const slides = document.querySelectorAll('.slide');
                const dots   = document.querySelectorAll('.dot');
                slides.forEach((s, idx) => s.classList.toggle('active', idx === {i}));
                dots.forEach((d, idx)   => d.classList.toggle('active', idx === {i}));
                document.documentElement.style.fontSize = {i} === 0 ? '130%' : '149.5%';
            """)
            await page.wait_for_timeout(300)
            path = Path(__file__).parent / f"_slide_{i+1:02d}.png"
            await page.screenshot(path=str(path), full_page=False)
            screenshots.append(path)
            print(f"  Screenshot {i+1}/{n_slides}")

        await browser.close()
    return screenshots


def build_pptx(screenshots):
    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W / 96)
    prs.slide_height = Inches(SLIDE_H / 96)

    blank_layout = prs.slide_layouts[6]

    for img_path in screenshots:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(img_path),
            left=0, top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    prs.save(str(OUT_PATH))
    print(f"\nSaved: {OUT_PATH}")


async def main():
    print("Capturing slides...")
    screenshots = await capture_slides()
    print("Building PPTX...")
    build_pptx(screenshots)
    for p in screenshots:
        p.unlink()
    print("Done!")


asyncio.run(main())
